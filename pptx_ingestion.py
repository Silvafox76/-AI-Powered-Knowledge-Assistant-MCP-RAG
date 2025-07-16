
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings

def ingest_pptx(file_path: str, collection_name: str = "my_documents"):
    """
    Ingests a PPTX file, extracts text, chunks it, generates embeddings, and stores them in ChromaDB.
    """
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_frame_text = shape.text_frame.text
                    if text_frame_text:
                        full_text.append(text_frame_text)
                if hasattr(shape, "table") and shape.table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                full_text.append(cell.text)

        text = "\n".join(full_text)

        # Chunk text
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            is_separator_regex=False,
        )
        chunks = text_splitter.create_documents([text])

        # Generate embeddings
        embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

        # Store in ChromaDB
        vectorstore = Chroma.from_documents(
            documents=chunks,
            embedding=embeddings,
            collection_name=collection_name,
            persist_directory="./chroma_db"
        )
        print(f"Successfully ingested {file_path} into ChromaDB collection {collection_name}")

    except Exception as e:
        print(f"Error ingesting PPTX {file_path}: {e}")

if __name__ == "__main__":
    import os
    os.makedirs("../data", exist_ok=True)
    # Example usage (requires a dummy.pptx in the data directory)
    # You would typically have your PPTX files here
    # For demonstration, let's assume you have a dummy.pptx in the data folder
    # ingest_pptx("../data/dummy.pptx")
    print("To test, create a PPTX file in the ../data directory and uncomment the ingest_pptx line.")


