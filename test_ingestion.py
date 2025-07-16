
import unittest
import os
from unittest.mock import patch, MagicMock

# Assuming these are in the parent directory of tests
from backend.ingestion.pdf_ingestion import ingest_pdf
from backend.ingestion.pptx_ingestion import ingest_pptx
from backend.ingestion.text_ingestion import ingest_text, ingest_docx
from backend.ingestion.csv_ingestion import ingest_csv
from backend.ingestion.web_ingestion import ingest_web_page

class TestIngestionModules(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        # Create a dummy data directory and files for testing
        cls.data_dir = "./test_data"
        os.makedirs(cls.data_dir, exist_ok=True)

        # Dummy PDF
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(os.path.join(cls.data_dir, "dummy.pdf"))
        c.drawString(100, 750, "This is a dummy PDF for testing.")
        c.save()

        # Dummy PPTX
        from pptx import Presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Dummy Presentation"
        prs.save(os.path.join(cls.data_dir, "dummy.pptx"))

        # Dummy TXT
        with open(os.path.join(cls.data_dir, "dummy.txt"), "w") as f:
            f.write("This is a dummy text file.")

        # Dummy DOCX
        from docx import Document
        document = Document()
        document.add_paragraph("This is a dummy DOCX file.")
        document.save(os.path.join(cls.data_dir, "dummy.docx"))

        # Dummy CSV
        import pandas as pd
        df = pd.DataFrame({"col1": [1, 2], "col2": ["a", "b"]})
        df.to_csv(os.path.join(cls.data_dir, "dummy.csv"), index=False)

    @classmethod
    def tearDownClass(cls):
        # Clean up dummy data directory and files
        import shutil
        shutil.rmtree(cls.data_dir)
        # Clean up chroma_db created by tests
        if os.path.exists("./chroma_db"):
            shutil.rmtree("./chroma_db")

    @patch("backend.ingestion.pdf_ingestion.Chroma.from_documents")
    @patch("backend.ingestion.pdf_ingestion.HuggingFaceEmbeddings")
    def test_ingest_pdf(self, mock_embeddings, mock_chroma):
        ingest_pdf(os.path.join(self.data_dir, "dummy.pdf"))
        mock_embeddings.assert_called_once()
        mock_chroma.assert_called_once()

    @patch("backend.ingestion.pptx_ingestion.Chroma.from_documents")
    @patch("backend.ingestion.pptx_ingestion.HuggingFaceEmbeddings")
    def test_ingest_pptx(self, mock_embeddings, mock_chroma):
        ingest_pptx(os.path.join(self.data_dir, "dummy.pptx"))
        mock_embeddings.assert_called_once()
        mock_chroma.assert_called_once()

    @patch("backend.ingestion.text_ingestion.Chroma.from_documents")
    @patch("backend.ingestion.text_ingestion.HuggingFaceEmbeddings")
    def test_ingest_text(self, mock_embeddings, mock_chroma):
        ingest_text(os.path.join(self.data_dir, "dummy.txt"))
        mock_embeddings.assert_called_once()
        mock_chroma.assert_called_once()

    @patch("backend.ingestion.text_ingestion.Chroma.from_documents")
    @patch("backend.ingestion.text_ingestion.HuggingFaceEmbeddings")
    def test_ingest_docx(self, mock_embeddings, mock_chroma):
        ingest_docx(os.path.join(self.data_dir, "dummy.docx"))
        mock_embeddings.assert_called_once()
        mock_chroma.assert_called_once()

    @patch("backend.ingestion.csv_ingestion.Chroma.from_documents")
    @patch("backend.ingestion.csv_ingestion.HuggingFaceEmbeddings")
    def test_ingest_csv(self, mock_embeddings, mock_chroma):
        ingest_csv(os.path.join(self.data_dir, "dummy.csv"))
        mock_embeddings.assert_called_once()
        mock_chroma.assert_called_once()

    @patch("backend.ingestion.web_ingestion.requests.get")
    @patch("backend.ingestion.web_ingestion.BeautifulSoup")
    @patch("backend.ingestion.web_ingestion.Chroma.from_documents")
    @patch("backend.ingestion.web_ingestion.HuggingFaceEmbeddings")
    def test_ingest_web_page(self, mock_embeddings, mock_chroma, mock_bs, mock_requests_get):
        mock_requests_get.return_value.raise_for_status.return_value = None
        mock_requests_get.return_value.text = "<html><body><p>Test web content.</p></body></html>"
        mock_bs.return_value.find_all.return_value = [MagicMock(get_text=lambda: "Test web content.")]

        ingest_web_page("http://example.com")
        mock_requests_get.assert_called_once_with("http://example.com")
        mock_embeddings.assert_called_once()
        mock_chroma.assert_called_once()

if __name__ == "__main__":
    unittest.main()


